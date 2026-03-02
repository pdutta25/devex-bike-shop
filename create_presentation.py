#!/usr/bin/env python3
"""Generate AI-Powered Developer Workflow demo presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colour palette ──────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x11, 0x17)   # near-black
BG_CARD      = RGBColor(0x16, 0x1B, 0x22)   # GitHub dark card
ACCENT_BLUE  = RGBColor(0x58, 0xA6, 0xFF)   # GitHub blue
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)   # GitHub green
ACCENT_AMBER = RGBColor(0xF5, 0xA6, 0x23)   # Amber/orange
ACCENT_PURPLE= RGBColor(0xA3, 0x71, 0xF7)   # Purple accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xC9, 0xD1, 0xD9)
MID_GRAY     = RGBColor(0x8B, 0x94, 0x9E)
DIM_GRAY     = RGBColor(0x48, 0x4F, 0x58)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────
def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(4)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=LIGHT_GRAY, spacing=Pt(8),
                    bullet_color=ACCENT_BLUE, icon_items=None):
    """Add a bulleted list. icon_items maps index -> emoji prefix."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = ""
        if icon_items and i in icon_items:
            prefix = icon_items[i] + "  "
        p.text = prefix + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_number_badge(slide, left, top, number, color=ACCENT_BLUE, size=Inches(0.55)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(0)
    return shape


def add_gradient_bar(slide, top, color1, color2, height=Pt(6)):
    """Fake gradient bar with two adjacent rectangles."""
    half = int(W / 2)
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, half, height)
    s1.fill.solid()
    s1.fill.fore_color.rgb = color1
    s1.line.fill.background()
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, half, top, half, height)
    s2.fill.solid()
    s2.fill.fore_color.rgb = color2
    s2.line.fill.background()


def lighter(color, dr=40, dg=20, db=20):
    return RGBColor(min(255, color[0] + dr), min(255, color[1] + dg), min(255, color[2] + db))


def make_section_slide(title, subtitle, act_label, accent_color=ACCENT_BLUE,
                       emoji=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    solid_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), accent_color, lighter(accent_color), Pt(6))
    # Act label pill
    pill = add_shape(slide, Inches(0.8), Inches(2.2), Inches(2.2), Inches(0.5),
                     accent_color, corner_radius=0.5)
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = act_label
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    # Title
    display_title = f"{emoji}  {title}" if emoji else title
    add_text(slide, Inches(0.8), Inches(3.0), Inches(10), Inches(1.2),
             display_title, 44, WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(4.3), Inches(9), Inches(0.8),
             subtitle, 22, MID_GRAY)
    return slide


def make_feature_slide(feature_num, title, subtitle, talking_point,
                       steps, accent_color=ACCENT_BLUE, emoji=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), accent_color,
                     lighter(accent_color), Pt(4))
    # Feature badge
    add_number_badge(slide, Inches(0.7), Inches(0.55), feature_num, accent_color)
    badge_label = f"FEATURE #{feature_num}"
    add_text(slide, Inches(1.4), Inches(0.55), Inches(3), Inches(0.45),
             badge_label, 14, accent_color, bold=True)
    # Title with emoji
    display_title = f"{emoji}  {title}" if emoji else title
    add_text(slide, Inches(0.7), Inches(1.2), Inches(11), Inches(0.8),
             display_title, 34, WHITE, bold=True)
    # Subtitle / nickname
    add_text(slide, Inches(0.7), Inches(2.0), Inches(11), Inches(0.5),
             subtitle, 18, MID_GRAY)
    add_accent_line(slide, Inches(0.7), Inches(2.6), Inches(1.5), accent_color)
    # Steps on left
    step_top = Inches(2.9)
    for i, step in enumerate(steps):
        add_text(slide, Inches(1.0), step_top + Inches(i * 0.50),
                 Inches(7.0), Inches(0.45), step, 15, LIGHT_GRAY)
        # Small bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(0.75), step_top + Inches(i * 0.50) + Pt(5),
            Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent_color
        dot.line.fill.background()

    # Talking point card on right
    card_left = Inches(8.5)
    card_top  = Inches(2.9)
    card_w    = Inches(4.3)
    card_h    = Inches(3.5)
    card = add_shape(slide, card_left, card_top, card_w, card_h, BG_CARD, 0.05)
    add_text(slide, card_left + Inches(0.3), card_top + Inches(0.25),
             card_w - Inches(0.6), Inches(0.4),
             "\U0001F4AC  SAY TO AUDIENCE", 12, accent_color, bold=True)
    add_text(slide, card_left + Inches(0.3), card_top + Inches(0.7),
             card_w - Inches(0.6), card_h - Inches(1.0),
             f"\u201C{talking_point}\u201D", 16, LIGHT_GRAY)
    return slide


# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, BG_DARK)
add_gradient_bar(slide, Inches(0), ACCENT_BLUE, ACCENT_PURPLE, Pt(6))

# Large centered title
add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
         "\U0001F3AC  AI-Powered Developer Workflow", 48, WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
         "From Idea to Production with GitHub Copilot", 28, ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)

# Decorative line
add_accent_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), ACCENT_BLUE)

# Subtitle info
add_text(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
         "Live Demo  \u2022  ~12 minutes  \u2022  10 Features", 20, MID_GRAY,
         alignment=PP_ALIGN.CENTER)

# Bottom badges
badge_labels = ["VS Code + Copilot", "GitHub Actions", "GitHub Copilot Chat",
                "AI Code Review"]
badge_colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_AMBER]
start_x = Inches(2.3)
for i, (label, clr) in enumerate(zip(badge_labels, badge_colors)):
    pill = add_shape(slide, start_x + Inches(i * 2.4), Inches(5.5),
                     Inches(2.1), Inches(0.45), clr, 0.5)
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, BG_DARK)
add_gradient_bar(slide, Inches(0), ACCENT_BLUE, ACCENT_PURPLE, Pt(4))
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.9),
         "Demo Agenda", 40, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2), ACCENT_BLUE)

acts = [
    ("ACT 1", "Write Code with Copilot", "5 min", ACCENT_BLUE,
     ["Code Completion", "Inline Chat", "Copilot Chat",
      "Context Awareness", "Test Generation"]),
    ("ACT 2", "Commit & Push", "2 min", ACCENT_GREEN,
     ["AI Commit Messages"]),
    ("ACT 3", "Create PR with AI", "2 min", ACCENT_PURPLE,
     ["AI PR Descriptions"]),
    ("ACT 4", "AI on GitHub.com", "3 min", ACCENT_AMBER,
     ["Code Review", "Repo Chat", "Actions Debugging"]),
]

card_w = Inches(2.8)
card_h = Inches(4.5)
gap = Inches(0.25)
total_w = 4 * card_w + 3 * gap
start_x = int((W - total_w) / 2)

for i, (act, title, duration, color, features) in enumerate(acts):
    cx = start_x + i * (card_w + gap)
    cy = Inches(2.0)
    card = add_shape(slide, cx, cy, card_w, card_h, BG_CARD, 0.04)
    # Act pill
    pill_w = Inches(1.2)
    pill = add_shape(slide, cx + Inches(0.15), cy + Inches(0.2),
                     pill_w, Inches(0.35), color, 0.5)
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.text = act
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    # Duration
    add_text(slide, cx + Inches(1.5), cy + Inches(0.22),
             Inches(1.1), Inches(0.3), duration, 12, MID_GRAY,
             alignment=PP_ALIGN.RIGHT)
    # Title
    add_text(slide, cx + Inches(0.2), cy + Inches(0.7),
             card_w - Inches(0.4), Inches(0.7), title, 18, WHITE, bold=True)
    # Features
    for j, feat in enumerate(features):
        add_text(slide, cx + Inches(0.45), cy + Inches(1.5) + Inches(j * 0.45),
                 card_w - Inches(0.6), Inches(0.4),
                 f"\u2022  {feat}", 14, LIGHT_GRAY)

# Bottom note
add_text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
         "Total demo time: ~12 minutes  \u2022  10 AI-powered features",
         16, MID_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — PRE-DEMO SETUP
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, BG_DARK)
add_gradient_bar(slide, Inches(0), MID_GRAY, DIM_GRAY, Pt(4))
add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.9),
         "\u2699\uFE0F  Pre-Demo Setup", 38, WHITE, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(2), MID_GRAY)

# Code block card
code_card = add_shape(slide, Inches(0.8), Inches(1.9), Inches(7), Inches(3.8),
                      RGBColor(0x0D, 0x11, 0x17), 0.03)
code_lines = [
    "# Ensure you're on main and up to date",
    "cd ~/Claude/devx-bike-shop",
    "git checkout main && git pull",
    "",
    "# Create a feature branch",
    "git checkout -b feature/spring-sale-banner",
    "",
    "# Open VS Code",
    "code .",
]
for i, line in enumerate(code_lines):
    clr = MID_GRAY if line.startswith("#") else ACCENT_GREEN
    if line == "":
        continue
    add_text(slide, Inches(1.1), Inches(2.1) + Inches(i * 0.4),
             Inches(6.5), Inches(0.35), line, 15, clr,
             font_name="Courier New")

# Checklist card
check_card = add_shape(slide, Inches(8.3), Inches(1.9), Inches(4.3),
                       Inches(3.8), BG_CARD, 0.04)
add_text(slide, Inches(8.6), Inches(2.1), Inches(3.8), Inches(0.4),
         "CHECKLIST", 14, ACCENT_BLUE, bold=True)
checks = [
    "\u2705  VS Code installed & open",
    "\u2705  Copilot extension active",
    "\u2705  GitHub CLI authenticated",
    "\u2705  Repo cloned & up to date",
    "\u2705  Feature branch created",
]
for i, item in enumerate(checks):
    add_text(slide, Inches(8.6), Inches(2.7) + Inches(i * 0.5),
             Inches(3.8), Inches(0.4), item, 15, LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════
# ACT 1 — Section Divider
# ════════════════════════════════════════════════════════════════════
make_section_slide(
    "Write Code with Copilot",
    "5 features  \u2022  5 minutes  \u2022  From blank file to tested component",
    "ACT 1", ACCENT_BLUE, "\u270D\uFE0F"
)

# ── Feature #1 ──
make_feature_slide(
    1, "Copilot Code Completion", '"The Mind Reader"',
    "I typed 2 lines of real code. Copilot understood my project and wrote the rest.",
    [
        "Create new file: src/components/ui/announcement-bar.tsx",
        'Type the component shell and export function stub',
        "STOP \u2014 let Copilot ghost-text appear",
        "Tab to accept each suggestion (countdown logic, state, JSX)",
        'Type natural comments if needed:',
        '  // calculate days, hours, minutes until March 31 2026',
        '  // don\'t show if user dismissed via localStorage',
    ],
    ACCENT_BLUE, "\U0001F9E0"
)

# ── Feature #2 ──
make_feature_slide(
    2, "Copilot Inline Chat", '"The Surgeon"',
    "I described what I wanted in English. It rewrote the code in place \u2014 no copy-paste, no Stack Overflow.",
    [
        "Select the entire component",
        "Press Cmd+I to open inline chat",
        "Type: style with gradient amber-500 \u2192 orange-600",
        "Add promo code SPRING26",
        "Add a Shop Now link to /bikes",
        "Add a dismiss X button",
        "Accept the inline edit",
    ],
    ACCENT_PURPLE, "\U0001FA7A"
)

# ── Feature #3 ──
make_feature_slide(
    3, "Copilot Chat Panel", '"The Pair Programmer"',
    "It knows the project structure. It\u2019s not just autocomplete \u2014 it understands architecture.",
    [
        "Open Copilot Chat panel (Cmd+Shift+I)",
        "Ask: @workspace where should I add AnnouncementBar",
        "  so it shows at the top of every storefront page?",
        "Copilot points to src/app/(storefront)/layout.tsx",
        "Click the suggested file to open it",
    ],
    ACCENT_GREEN, "\U0001F91D"
)

# ── Feature #4 ──
make_feature_slide(
    4, "Auto-Import & Context Awareness", '"Context Awareness"',
    "It auto-imported from the right path. It knows where my component lives.",
    [
        "Open src/app/(storefront)/layout.tsx",
        "Find where {children} are rendered",
        "Above {children}, start typing <Announce...",
        "Copilot suggests <AnnouncementBar />",
        "AND auto-adds the import at the top of the file",
    ],
    ACCENT_AMBER, "\U0001F50D"
)

# ── Feature #5 ──
make_feature_slide(
    5, "Copilot Generates Tests", '"The QA Engineer"',
    "It wrote tests covering edge cases I didn\u2019t even think of. Including localStorage mocking.",
    [
        "Open Copilot Chat \u2192 type: /tests for AnnouncementBar",
        "Copilot generates test cases:",
        "  \u2022 Renders the banner",
        "  \u2022 Countdown displays correctly",
        "  \u2022 Dismiss button hides it",
        "  \u2022 localStorage persistence",
        "Save to __tests__/announcement-bar.test.tsx",
    ],
    RGBColor(0xE0, 0x55, 0x55), "\U0001F9EA"
)


# ════════════════════════════════════════════════════════════════════
# ACT 2 — Section Divider
# ════════════════════════════════════════════════════════════════════
make_section_slide(
    "Commit & Push",
    "AI-generated commit messages that actually make sense",
    "ACT 2", ACCENT_GREEN, "\U0001F4E6"
)

# ── Feature #6 ──
make_feature_slide(
    6, "Copilot Commit Message", '"No More \'fixed stuff\'"',
    "It analyzed the diff and wrote a conventional commit message. No more \u2018updated stuff\u2019.",
    [
        "Go to Source Control panel in VS Code",
        "Stage the changed files",
        "Click the \u2728 sparkle icon next to commit message box",
        "Copilot generates:",
        "  feat: add spring sale announcement bar",
        "       with countdown timer",
        "Commit and push to feature/spring-sale-banner",
    ],
    ACCENT_GREEN, "\u2728"
)


# ════════════════════════════════════════════════════════════════════
# ACT 3 — Section Divider
# ════════════════════════════════════════════════════════════════════
make_section_slide(
    "Create PR with AI",
    "From diff to documentation in one click",
    "ACT 3", ACCENT_PURPLE, "\U0001F4DD"
)

# ── Feature #7 ──
make_feature_slide(
    7, "Copilot PR Description", '"The Technical Writer"',
    "From diff to documentation in one click. This is what your PR descriptions should look like.",
    [
        "Go to GitHub \u2192 see 'Compare & pull request' banner",
        "Click it to create a new PR",
        "In description box, click the Copilot \u2728 button",
        "Copilot generates full PR description:",
        "  \u2022 Summary of changes",
        "  \u2022 Files modified",
        "  \u2022 What to test",
        "  \u2022 Screenshots suggestion",
        "Create the PR \u2192 triggers Actions workflow",
    ],
    ACCENT_PURPLE, "\U0001F4CB"
)


# ════════════════════════════════════════════════════════════════════
# ACT 4 — Section Divider
# ════════════════════════════════════════════════════════════════════
make_section_slide(
    "AI on GitHub.com",
    "Code review, repo Q&A, and CI debugging \u2014 all AI-powered",
    "ACT 4", ACCENT_AMBER, "\U0001F310"
)

# ── Feature #8 ──
make_feature_slide(
    8, "Copilot Code Review", '"The Senior Reviewer"',
    "A senior engineer review in 30 seconds. It catches things humans miss at 2pm on a Friday.",
    [
        "On the PR page \u2192 click Copilot \u2192 Review",
        "(or it auto-reviews)",
        "Copilot leaves inline comments like:",
        '  \u2022 "Consider memoizing the countdown interval"',
        '  \u2022 "The localStorage key should be a constant"',
        "  \u2022 Potential accessibility suggestions",
    ],
    ACCENT_AMBER, "\U0001F50E"
)

# ── Feature #9 ──
make_feature_slide(
    9, "Copilot Chat on GitHub.com", '"Ask Your Repo Anything"',
    "New team member joins? They don\u2019t need to read 50 files. They ask the repo directly.",
    [
        "Click the Copilot icon in the GitHub header",
        "Ask: What does the checkout flow look like?",
        "Ask: What API endpoints are available?",
        "Ask: Summarize what this PR changes",
    ],
    RGBColor(0x79, 0xC0, 0xFF), "\U0001F4AC"
)

# ── Feature #10 ──
make_feature_slide(
    10, "Copilot in Actions", '"The DevOps Whisperer"',
    "No more googling cryptic build errors. AI explains the failure and suggests the fix.",
    [
        "Click on the Actions tab \u2192 click the running workflow",
        "If any step fails \u2192 Copilot shows 'Fix with Copilot'",
        "Or in Copilot Chat:",
        "  @github why did my last workflow run fail?",
        "Copilot explains the error and suggests a fix",
    ],
    RGBColor(0xE0, 0x55, 0x55), "\U0001F6E0\uFE0F"
)


# ════════════════════════════════════════════════════════════════════
# FINALE SLIDE
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, BG_DARK)
add_gradient_bar(slide, Inches(0), ACCENT_BLUE, ACCENT_PURPLE, Pt(6))

add_text(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0),
         "\U0001F3AF  The Punchline", 42, WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

recap_items = [
    ("\u2705", "Copilot wrote the component",           ACCENT_BLUE),
    ("\U0001F3A8", "Copilot styled it",                  ACCENT_PURPLE),
    ("\U0001F4CD", "Copilot told me where to put it",    ACCENT_GREEN),
    ("\U0001F9EA", "Copilot wrote the tests",            RGBColor(0xE0, 0x55, 0x55)),
    ("\U0001F4DD", "Copilot wrote the commit message",   ACCENT_GREEN),
    ("\U0001F4CB", "Copilot wrote the PR description",   ACCENT_PURPLE),
    ("\U0001F50E", "Copilot reviewed the code",          ACCENT_AMBER),
    ("\U0001F6E0\uFE0F", "Copilot explained the build",  RGBColor(0xE0, 0x55, 0x55)),
]

# Two columns of recap
col1_x = Inches(1.5)
col2_x = Inches(7.0)
start_y = Inches(1.6)
row_h = Inches(0.52)

for i, (emoji, text, color) in enumerate(recap_items):
    col_x = col1_x if i < 4 else col2_x
    row_y = start_y + (i % 4) * row_h
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        col_x, row_y + Pt(4), Pt(12), Pt(12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, col_x + Inches(0.25), row_y,
             Inches(4.5), Inches(0.45),
             f"{emoji}   {text}", 18, LIGHT_GRAY)

# Highlight box
highlight = add_shape(slide, Inches(2.5), Inches(4.1), Inches(8.3), Inches(1.1),
                      BG_CARD, 0.05)
add_text(slide, Inches(2.8), Inches(4.2), Inches(7.7), Inches(0.9),
         "I typed maybe 20 words of actual code. The rest was AI.", 24,
         ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom statement
add_text(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(1.0),
         "This is how developers work in 2026.", 36, WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

# Footer accent
add_gradient_bar(slide, Inches(7.2), ACCENT_BLUE, ACCENT_PURPLE, Pt(4))


# ════════════════════════════════════════════════════════════════════
# SLIDE — THANK YOU
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, BG_DARK)
add_gradient_bar(slide, Inches(0), ACCENT_BLUE, ACCENT_PURPLE, Pt(6))

add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
         "Thank You!", 54, WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_accent_line(slide, Inches(5.5), Inches(3.3), Inches(2.3), ACCENT_BLUE)
add_text(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.7),
         "Questions?", 28, ACCENT_BLUE,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
         "github.com/copilot  \u2022  code.visualstudio.com", 18, MID_GRAY,
         alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════
output_path = "/Users/pinakd/Claude/devx-bike-shop/AI-Powered-Developer-Workflow-Demo.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
